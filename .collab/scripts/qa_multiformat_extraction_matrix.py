#!/usr/bin/env python3
"""Exercise DocTalk's non-PDF extraction pipeline with generated fixtures.

The fixtures are generated in memory so this does not mutate `test_inputs/`.
This matrix covers format-specific extraction, markdown table preservation,
section title detection, chunking compatibility, and upload-content validation.
"""

from __future__ import annotations

import argparse
import io
import json
import sys
import zipfile
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

ROOT = Path(__file__).resolve().parents[2]
BACKEND_DIR = ROOT / "backend"
if str(BACKEND_DIR) not in sys.path:
    sys.path.insert(0, str(BACKEND_DIR))


@dataclass(frozen=True)
class FormatCase:
    name: str
    file_type: str
    filename: str
    content_type: str
    data: bytes
    expected_terms: tuple[str, ...]
    expected_title: str | None = None
    expect_table: bool = False
    min_pages: int = 1
    min_chunks: int = 1


@dataclass(frozen=True)
class InvalidCase:
    name: str
    file_type: str
    data: bytes
    expected_valid: bool


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser()
    parser.add_argument("--json-out", required=True)
    return parser.parse_args()


def make_docx() -> bytes:
    from docx import Document

    doc = Document()
    doc.add_heading("Quarterly Risk Memo", level=1)
    doc.add_paragraph("DocTalk should extract DOCX paragraphs and cite the risk memo text.")
    table = doc.add_table(rows=3, cols=2)
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Retention"
    table.cell(1, 1).text = "94%"
    table.cell(2, 0).text = "Escalations"
    table.cell(2, 1).text = "2"
    doc.add_heading("Customer Notes", level=2)
    doc.add_paragraph("Chinese coverage: 客户反馈稳定。")
    out = io.BytesIO()
    doc.save(out)
    return out.getvalue()


def make_pptx() -> bytes:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "Launch Readiness"
    box = slide.shapes.add_textbox(Inches(0.7), Inches(1.3), Inches(8), Inches(1))
    box.text_frame.text = "The launch plan requires support coverage and rollback checks."
    table = slide.shapes.add_table(3, 2, Inches(0.7), Inches(2.4), Inches(6), Inches(1.2)).table
    table.cell(0, 0).text = "Owner"
    table.cell(0, 1).text = "Status"
    table.cell(1, 0).text = "Support"
    table.cell(1, 1).text = "Ready"
    table.cell(2, 0).text = "Rollback"
    table.cell(2, 1).text = "Documented"
    try:
        notes = slide.notes_slide.notes_text_frame
        notes.text = "Speaker notes include the escalation channel."
    except Exception:
        pass
    out = io.BytesIO()
    prs.save(out)
    return out.getvalue()


def make_xlsx() -> bytes:
    from openpyxl import Workbook

    wb = Workbook()
    ws = wb.active
    ws.title = "Revenue"
    ws.append(["Quarter", "Revenue", "Region"])
    ws.append(["Q1", 1200, "NA"])
    ws.append(["Q2", 1450, "EU"])
    ws2 = wb.create_sheet("中文")
    ws2.append(["主题", "结论"])
    ws2.append(["资金流", "净流入"])
    out = io.BytesIO()
    wb.save(out)
    wb.close()
    return out.getvalue()


def make_zip_without_content_types() -> bytes:
    out = io.BytesIO()
    with zipfile.ZipFile(out, "w") as zf:
        zf.writestr("word/document.xml", "<document />")
    return out.getvalue()


def cases() -> list[FormatCase]:
    return [
        FormatCase(
            name="docx_paragraphs_table_cjk",
            file_type="docx",
            filename="qa-risk-memo.docx",
            content_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            data=make_docx(),
            expected_terms=("risk memo", "Retention", "客户反馈"),
            expected_title="Quarterly Risk Memo",
            expect_table=True,
        ),
        FormatCase(
            name="pptx_slide_table_notes",
            file_type="pptx",
            filename="qa-launch-readiness.pptx",
            content_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            data=make_pptx(),
            expected_terms=("launch plan", "Rollback", "Speaker Notes"),
            expected_title="Launch Readiness",
            expect_table=True,
        ),
        FormatCase(
            name="xlsx_multisheet_markdown_tables",
            file_type="xlsx",
            filename="qa-revenue.xlsx",
            content_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            data=make_xlsx(),
            expected_terms=("Quarter", "Revenue", "资金流"),
            expected_title="Revenue",
            expect_table=True,
            min_pages=2,
        ),
        FormatCase(
            name="txt_utf8_long_en_cjk",
            file_type="txt",
            filename="qa-notes.txt",
            content_type="text/plain",
            data=("Plain text extraction should preserve user notes.\n\n中文段落用于验证 UTF-8 文本。\n" * 8).encode(),
            expected_terms=("Plain text extraction", "中文段落"),
        ),
        FormatCase(
            name="md_headings_and_table",
            file_type="md",
            filename="qa-plan.md",
            content_type="text/markdown",
            data=(
                "# Product Plan\n\n"
                "Markdown extraction should keep paragraphs and tables.\n\n"
                "| Area | Priority |\n| --- | --- |\n| Citations | High |\n| URL import | Medium |\n\n"
                "## 中文章节\n\n引用跳转需要稳定。\n"
            ).encode(),
            expected_terms=("Markdown extraction", "Citations", "引用跳转"),
            expected_title="Product Plan",
            expect_table=True,
        ),
    ]


def invalid_cases() -> list[InvalidCase]:
    return [
        InvalidCase("invalid_docx_magic_bytes", "docx", b"not a zip", False),
        InvalidCase("docx_zip_without_content_types", "docx", make_zip_without_content_types(), False),
        InvalidCase("invalid_pptx_magic_bytes", "pptx", b"PK\x03\x04broken", False),
        InvalidCase("txt_without_magic_is_allowed", "txt", b"hello", True),
        InvalidCase("md_without_magic_is_allowed", "md", b"# hello", True),
    ]


def pages_to_chunkable(extracted: list[Any]) -> list[Any]:
    from app.services.parse_service import BlockInfo, PageInfo

    pages = []
    for page in extracted:
        blocks = []
        if page.section_title:
            blocks.append(
                BlockInfo(
                    page=page.page_number,
                    text=page.section_title,
                    bbox=(0.0, 0.0, 1.0, 0.08),
                    font_size=18.0,
                )
            )
        blocks.append(
            BlockInfo(
                page=page.page_number,
                text=page.text,
                bbox=(0.0, 0.1 if page.section_title else 0.0, 1.0, 1.0),
                font_size=12.0,
            )
        )
        pages.append(
            PageInfo(
                page_number=page.page_number,
                width_pt=page.width_pt or 612.0,
                height_pt=page.height_pt or 792.0,
                rotation=0,
                blocks=blocks,
            )
        )
    return pages


def run_case(case: FormatCase) -> dict[str, Any]:
    from app.api.documents import _validate_file_content
    from app.services.extractors.base import extract_document
    from app.services.parse_service import ParseService
    from app.services.table_service import parse_markdown_tables

    extracted = extract_document(case.data, case.file_type)
    text = "\n\n".join(page.text for page in extracted)
    titles = [page.section_title for page in extracted if page.section_title]
    chunkable_pages = pages_to_chunkable(extracted)
    chunks = ParseService().chunk_document(chunkable_pages)
    chunk_text = "\n\n".join(chunk.text for chunk in chunks)
    tables = parse_markdown_tables(text)

    checks = [
        ("upload_content_validation_accepts_fixture", _validate_file_content(case.data, case.file_type)),
        ("page_count", len(extracted) >= case.min_pages),
        ("chunk_count", len(chunks) >= case.min_chunks),
        ("expected_terms_in_extracted_text", all(term in text for term in case.expected_terms)),
        ("expected_terms_in_chunks", all(term in chunk_text for term in case.expected_terms)),
    ]
    if case.expected_title:
        checks.append(("section_title_detected", case.expected_title in titles))
    if case.expect_table:
        checks.append(("markdown_table_detected", len(tables) > 0))

    return {
        "name": case.name,
        "file_type": case.file_type,
        "filename": case.filename,
        "content_type": case.content_type,
        "bytes": len(case.data),
        "page_count": len(extracted),
        "section_titles": titles[:5],
        "chunk_count": len(chunks),
        "table_count": len(tables),
        "first_page_preview": text[:500],
        "checks": [{"name": name, "result": "pass" if ok else "fail"} for name, ok in checks],
        "result": "pass" if all(ok for _, ok in checks) else "fail",
    }


def run_invalid_case(case: InvalidCase) -> dict[str, Any]:
    from app.api.documents import _validate_file_content

    actual = _validate_file_content(case.data, case.file_type)
    return {
        "name": case.name,
        "file_type": case.file_type,
        "bytes": len(case.data),
        "expected_valid": case.expected_valid,
        "actual_valid": actual,
        "result": "pass" if actual is case.expected_valid else "fail",
    }


def main() -> int:
    args = parse_args()
    report = {
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "scope": "non-pdf extractor/chunking/upload-validation matrix",
        "format_cases": [],
        "invalid_cases": [],
    }
    for case in cases():
        result = run_case(case)
        report["format_cases"].append(result)
        print(f"CASE {case.name}: {result['result'].upper()} pages={result['page_count']} chunks={result['chunk_count']}")
    for case in invalid_cases():
        result = run_invalid_case(case)
        report["invalid_cases"].append(result)
        print(f"INVALID {case.name}: {result['result'].upper()} actual_valid={result['actual_valid']}")

    all_results = report["format_cases"] + report["invalid_cases"]
    failures = [item["name"] for item in all_results if item["result"] != "pass"]
    report["summary"] = {
        "format_total": len(report["format_cases"]),
        "invalid_total": len(report["invalid_cases"]),
        "failed": failures,
    }
    report["result"] = "fail" if failures else "pass"

    out = Path(args.json_out)
    out.parent.mkdir(parents=True, exist_ok=True)
    out.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(f"MULTIFORMAT_EXTRACTION {report['result'].upper()}: wrote {args.json_out}")
    return 0 if not failures else 1


if __name__ == "__main__":
    raise SystemExit(main())
