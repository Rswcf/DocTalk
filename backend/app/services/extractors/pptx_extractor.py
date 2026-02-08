from __future__ import annotations

import io
from typing import List

from .base import ExtractedPage


def _table_to_markdown(table) -> str:
    """Convert a python-pptx Table to a markdown table string."""
    rows: List[List[str]] = []
    for row in table.rows:
        cells = [cell.text.strip().replace('\n', ' ') for cell in row.cells]
        rows.append(cells)

    if not rows:
        return ''

    lines: List[str] = []
    lines.append('| ' + ' | '.join(rows[0]) + ' |')
    lines.append('| ' + ' | '.join('---' for _ in rows[0]) + ' |')
    for row in rows[1:]:
        while len(row) < len(rows[0]):
            row.append('')
        lines.append('| ' + ' | '.join(row[:len(rows[0])]) + ' |')

    return '\n'.join(lines)


def extract_pptx(file_bytes: bytes) -> List[ExtractedPage]:
    """Extract text from PPTX files. Each slide = one page.

    Extracts text frames, tables (as markdown), and speaker notes.
    """
    from pptx import Presentation

    prs = Presentation(io.BytesIO(file_bytes))
    pages: List[ExtractedPage] = []

    for i, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        title = None

        for shape in slide.shapes:
            if shape.has_text_frame:
                shape_text = shape.text_frame.text.strip()
                if shape_text:
                    texts.append(shape_text)
            if hasattr(shape, 'text') and shape == slide.shapes.title:
                title = shape.text.strip()[:200] if shape.text else None

            # Handle tables in slides
            if shape.has_table:
                md_table = _table_to_markdown(shape.table)
                if md_table:
                    texts.append('')
                    texts.append(md_table)
                    texts.append('')

        # Extract speaker notes
        try:
            if slide.has_notes_slide and slide.notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    texts.append('')
                    texts.append('**Speaker Notes:**')
                    texts.append(notes_text)
        except Exception:
            pass  # notes_slide can raise if notes part is malformed

        text = '\n'.join(texts)
        if text.strip():
            pages.append(ExtractedPage(
                page_number=i,
                text=text,
                section_title=title,
            ))

    if not pages:
        pages.append(ExtractedPage(page_number=1, text='(empty presentation)'))

    return pages
